"""Upload routes — PDF/DOCX/TXT ingestion with full pipeline."""
from fastapi import APIRouter, Depends, UploadFile, File, Request
from src.core.security import get_current_user
from src.core.limiter import limiter
from src.core.constants import RATE_LIMIT_UPLOAD
from src.utils.validators import validate_file
from src.rag.ingestion.pdf_loader import load_pdf
from src.rag.ingestion.docx_loader import load_docx
from src.rag.ingestion.text_loader import load_text
from src.rag.ingestion.embedding_pipeline import run_ingestion_pipeline
from src.rag.vectorstores.pinecone_store import get_vector_store
from src.database.mongodb.repositories.document_repository import save_document_metadata
from src.core.logger import get_logger

router = APIRouter()
logger = get_logger(__name__)


@router.post("/upload")
@limiter.limit(RATE_LIMIT_UPLOAD)
async def upload_document(
    request: Request,
    file: UploadFile = File(...),
    current_user=Depends(get_current_user),
):
    user_id = str(current_user["_id"])
    filename = file.filename

    # Validate before reading to prevent OOM
    validate_file(filename, file.size or 0)

    content = await file.read()

    # Load based on extension
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    # Image types — store as binary passthrough document (metadata + filename)
    IMAGE_EXTS = {"png", "jpg", "jpeg", "gif", "webp", "bmp", "tiff", "svg"}
    # Spreadsheet types
    SHEET_EXTS = {"xlsx", "xls", "csv"}
    # Presentation types
    PRES_EXTS = {"pptx", "ppt"}
    # Text-like types (handled by text loader)
    TEXT_EXTS = {"txt", "md", "mdx", "rst", "html", "htm", "json", "jsonl",
                 "xml", "yaml", "yml", "csv",
                 "py", "js", "ts", "jsx", "tsx", "java", "cpp", "c", "cs",
                 "go", "rs", "rb", "php", "sh", "sql", "rtf", "odt", "doc"}

    if ext == "pdf":
        docs = await load_pdf(content, filename)
    elif ext == "docx":
        docs = await load_docx(content, filename)
    elif ext in IMAGE_EXTS:
        # Store image as a document with binary size metadata; text content is a placeholder
        from langchain_core.documents import Document
        import base64
        docs = [Document(
            page_content=f"[Image file: {filename} ({len(content) // 1024} KB)]",
            metadata={"filename": filename, "source": filename, "file_type": ext, "page": 0,
                      "data_uri": "data:image/" + ext + ";base64," + base64.b64encode(content[:4096]).decode()},
        )]
    elif ext in SHEET_EXTS and ext != "csv":
        # Excel — use openpyxl to extract cell text
        try:
            import io
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            rows = []
            for ws in wb.worksheets:
                rows.append(f"=== Sheet: {ws.title} ===")
                for row in ws.iter_rows(values_only=True):
                    rows.append("\t".join(str(c) if c is not None else "" for c in row))
            docs = await load_text("\n".join(rows).encode(), filename)
        except Exception as xlsx_err:
            logger.warning(f"Excel parse failed ({xlsx_err}), falling back to text loader")
            docs = await load_text(content, filename)
    elif ext in PRES_EXTS:
        # PowerPoint — extract slide text via python-pptx
        try:
            import io
            from pptx import Presentation as PptxPresentation
            prs = PptxPresentation(io.BytesIO(content))
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
                slides_text.append(f"=== Slide {i} ===\n" + "\n".join(texts))
            docs = await load_text("\n".join(slides_text).encode(), filename)
        except Exception as pptx_err:
            logger.warning(f"PPTX parse failed ({pptx_err}), falling back to text loader")
            docs = await load_text(content, filename)
    else:
        # Covers: txt, md, json, yaml, html, csv, code files, etc.
        docs = await load_text(content, filename)

    # Run ingestion pipeline (chunk → enrich → index)
    vector_store = get_vector_store()
    chunk_count = await run_ingestion_pipeline(
        documents=docs,
        user_id=user_id,
        filename=filename,
        vector_store=vector_store,
        use_semantic_chunking=True,
    )

    # Persist metadata to MongoDB
    await save_document_metadata(
        user_id=user_id,
        filename=filename,
        file_type=ext,
        chunk_count=chunk_count,
    )

    return {
        "message": "Document ingested and indexed successfully",
        "filename": filename,
        "chunks": chunk_count,
    }


from fastapi.responses import StreamingResponse
import json
import asyncio

@router.post("/upload/stream")
@limiter.limit(RATE_LIMIT_UPLOAD)
async def upload_document_stream(
    request: Request,
    file: UploadFile = File(...),
    current_user=Depends(get_current_user),
):
    user_id = str(current_user["_id"])
    filename = file.filename
    content = await file.read()

    async def event_generator():
        try:
            # Step 1: Parsing
            yield f"data: {json.dumps({'stage': 'parsing', 'progress': 10})}\n\n"
            validate_file(filename, len(content))
            ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

            IMAGE_EXTS = {"png", "jpg", "jpeg", "gif", "webp", "bmp", "tiff", "svg"}
            SHEET_EXTS = {"xlsx", "xls"}
            PRES_EXTS = {"pptx", "ppt"}

            if ext == "pdf":
                docs = await load_pdf(content, filename)
            elif ext == "docx":
                docs = await load_docx(content, filename)
            elif ext in IMAGE_EXTS:
                from langchain_core.documents import Document
                import base64
                docs = [Document(
                    page_content=f"[Image file: {filename} ({len(content) // 1024} KB)]",
                    metadata={"filename": filename, "source": filename, "file_type": ext, "page": 0},
                )]
            elif ext in SHEET_EXTS:
                try:
                    import io
                    import openpyxl
                    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
                    rows = []
                    for ws in wb.worksheets:
                        rows.append(f"=== Sheet: {ws.title} ===")
                        for row in ws.iter_rows(values_only=True):
                            rows.append("\t".join(str(c) if c is not None else "" for c in row))
                    docs = await load_text("\n".join(rows).encode(), filename)
                except Exception as xlsx_err:
                    logger.warning(f"Excel parse failed ({xlsx_err}), falling back to text loader")
                    docs = await load_text(content, filename)
            elif ext in PRES_EXTS:
                try:
                    import io
                    from pptx import Presentation as PptxPresentation
                    prs = PptxPresentation(io.BytesIO(content))
                    slides_text = []
                    for i, slide in enumerate(prs.slides, 1):
                        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
                        slides_text.append(f"=== Slide {i} ===\n" + "\n".join(texts))
                    docs = await load_text("\n".join(slides_text).encode(), filename)
                except Exception as pptx_err:
                    logger.warning(f"PPTX parse failed ({pptx_err}), falling back to text loader")
                    docs = await load_text(content, filename)
            else:
                docs = await load_text(content, filename)
            
            # Step 2: Chunking (simulate part of pipeline)
            yield f"data: {json.dumps({'stage': 'chunking', 'progress': 30})}\n\n"
            # We call the full pipeline which handles chunking and indexing, 
            # but we can emit progress before and after.
            from src.rag.ingestion.embedding_pipeline import clean_documents, semantic_chunk, enrich_metadata
            from src.rag.vectorstores.pinecone_store import get_vector_store

            vector_store = get_vector_store()
            
            # Clean
            docs = clean_documents(docs)
            
            # Chunk
            chunks = semantic_chunk(docs)
            yield f"data: {json.dumps({'stage': 'embedding', 'progress': 60, 'chunks': len(chunks)})}\n\n"
            
            # Enrich & Index
            chunks = enrich_metadata(chunks, user_id=user_id, filename=filename)
            
            yield f"data: {json.dumps({'stage': 'indexing', 'progress': 80})}\n\n"
            # In real-world, embedding and indexing happens here in one go with Pinecone
            await asyncio.to_thread(vector_store.add_documents, chunks)
            chunk_count = len(chunks)

            # Persist metadata to MongoDB
            await save_document_metadata(
                user_id=user_id,
                filename=filename,
                file_type=ext,
                chunk_count=chunk_count,
            )

            # Done
            yield f"data: {json.dumps({'stage': 'done', 'progress': 100, 'chunks': chunk_count})}\n\n"

        except Exception as e:
            logger.error(f"Streaming upload failed: {e}")
            yield f"data: {json.dumps({'error': str(e)})}\n\n"

    return StreamingResponse(event_generator(), media_type="text/event-stream")


@router.get("")
async def list_documents(current_user=Depends(get_current_user)):
    from src.database.mongodb.repositories.document_repository import get_user_documents
    user_id = str(current_user["_id"])
    return await get_user_documents(user_id)


@router.delete("/{document_id}")
async def delete_document(document_id: str, filename: str, current_user=Depends(get_current_user)):
    from src.database.mongodb.repositories.document_repository import delete_document_metadata
    from src.rag.vectorstores.pinecone_store import delete_by_filename
    
    user_id = str(current_user["_id"])
    
    # Delete from MongoDB
    db_success = await delete_document_metadata(user_id, document_id)
    
    # Delete from Pinecone
    pinecone_success = delete_by_filename(user_id, filename)
    
    if db_success or pinecone_success:
        return {"message": "Document deleted"}
    
    from fastapi import HTTPException
    raise HTTPException(status_code=404, detail="Document not found or could not be deleted")
